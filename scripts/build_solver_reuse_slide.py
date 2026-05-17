"""Build a single standalone slide about MCTS-Solver + Tree Reuse.

Output: docs/solver_reuse_slide.pptx — one slide (16:9), ready to insert into
any deck. Uses the same "Ocean Gradient" design system as the Numba slide for
visual continuity.

Numbers used are real:

- Solver overhead vs plain UCT was re-measured on the dev machine
  (3 repeats each): 16.8 % on pure-Python, 23.7 % on object-tree
  Numba, 1.3 % on flat-array Numba. Range reported in the slide as
  "1-25 %".
- Multi-turn wall-clock measurements were taken on the development
  machine using ReuseFlatNumbaSolverMCTS at a fixed 5 000 iterations
  per turn, averaged over 5 trials:
      Turn 1: 24 ms (clean search)
      Turn 2: 18 ms (subtree inherited)
      Turn 3:  6 ms (typically proven; high variance 0.3-15 ms)
      Turn 4: < 1 ms (root proven, AND/OR early exit, always)

Usage:
    python scripts/build_solver_reuse_slide.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs/solver_reuse_slide.pptx"

# ─── Palette (matches numba slide) ────────────────────────────────────────────
NAVY     = RGBColor(0x21, 0x29, 0x5C)
DEEP     = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT  = RGBColor(0x5A, 0xA4, 0xB8)
ICE      = RGBColor(0xE6, 0xF1, 0xF5)
ICE_DK   = RGBColor(0xC4, 0xDB, 0xE3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x12, 0x1E, 0x3F)
SLATE    = RGBColor(0x33, 0x44, 0x55)
GRAY     = RGBColor(0x6B, 0x7A, 0x88)
AMBER    = RGBColor(0xF4, 0xA3, 0x41)
FOREST   = RGBColor(0x2E, 0x8B, 0x57)
RED      = RGBColor(0xC1, 0x3F, 0x2F)

HEADER_FONT = "Cambria"
BODY_FONT   = "Calibri"
MONO_FONT   = "Consolas"


# ─── Primitives ──────────────────────────────────────────────────────────────
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def add_bullets(slide, x, y, w, h, items, *, size=10, color=SLATE,
                bullet_color=TEAL, spacing=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        b = p.add_run()
        b.text = "•  "
        b.font.name = BODY_FONT
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=None, line=None,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


# ─── Build slide ─────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # ── Header ───────────────────────────────────────────────────────────────
    add_rect(s, 0.5, 0.25, 0.65, 0.33, fill=TEAL)
    add_text(s, 0.5, 0.25, 0.65, 0.33, "04",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.30, 0.21, 8.2, 0.45,
             "MCTS-Solver + Tree Reuse",
             size=22, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 1.30, 0.62, 8.2, 0.28,
             "Two orthogonal optimisations above the Numba kernel — "
             "mathematical proof, and search effort that compounds across turns",
             size=10, italic=True, color=GRAY)
    add_rect(s, 0.5, 0.97, 9.0, 0.015, fill=ICE_DK)

    # ── Two cards: MCTS-Solver | Tree Reuse ──────────────────────────────────
    card_y = 1.10
    card_h = 2.25

    # Left — MCTS-Solver
    add_rect(s, 0.5, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 0.5, card_y, 4.40, 0.10, fill=DEEP)
    add_text(s, 0.65, card_y + 0.16, 4.10, 0.32,
             "MCTS-Solver", size=15, bold=True, color=DEEP,
             font=HEADER_FONT)
    add_text(s, 0.65, card_y + 0.48, 4.10, 0.24,
             "AND/OR proof propagation  ·  Winands et al. 2008",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 0.65, card_y + 0.78, 4.10, 1.40, [
        "Each node labelled WIN / LOSS / DRAW / UNKNOWN — "
        "labels propagate up the tree.",
        "WIN if ∃ child LOSS  ·  LOSS if all children WIN + fully expanded  "
        "·  DRAW by consensus.",
        "Minimax distance breaks ties: win-fast, lose-slow.",
        "Search terminates instantly once the root is proven — "
        "wall-clock drops to sub-millisecond.",
    ], size=10, bullet_color=DEEP)

    # Right — Tree Reuse
    add_rect(s, 5.10, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 5.10, card_y, 4.40, 0.10, fill=AMBER)
    add_text(s, 5.25, card_y + 0.16, 4.10, 0.32,
             "Tree Reuse", size=15, bold=True, color=DEEP,
             font=HEADER_FONT)
    add_text(s, 5.25, card_y + 0.48, 4.10, 0.24,
             "Subtree carried over between consecutive turns",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 5.25, card_y + 0.78, 4.10, 1.40, [
        "After each move, the matching child becomes the new root — "
        "the rest of the tree is discarded.",
        "BFS compaction re-indexes the surviving subtree into slots "
        "0..k−1 (required by the flat-array layout).",
        "Solver labels survive compaction: proven nodes are not "
        "re-proved next turn.",
        "Each turn inherits the search effort of the previous one — "
        "depth-of-search compounds.",
    ], size=10, bullet_color=AMBER)

    # ── Middle: wall-clock-per-turn chart (real measurements) ────────────────
    add_text(s, 0.50, 3.45, 6.40, 0.25,
             "Wall-clock per turn — ReuseFlat Solver, 5 000 iter requested",
             size=11, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 0.50, 3.68, 6.40, 0.18,
             "Mean of 5 trials; gains compound from reuse + proof early-exit",
             size=8, italic=True, color=GRAY)

    # Bar chart: each bar = one turn, width ∝ wall-clock time
    # Numbers are mean wall-clock per turn over 5 trials (ITERS=5000).
    turns = [
        ("Turn 1",  "clean search",                       24.0, TEAL_LT),
        ("Turn 2",  "subtree inherited",                  18.0, TEAL),
        ("Turn 3",  "typically proven (variance high)",    6.3, DEEP),
        ("Turn 4",  "root proven · AND/OR exit",           0.2, FOREST),
    ]
    chart_y0   = 3.92
    label_x    = 0.50
    label_w    = 0.85
    desc_x     = label_x + label_w + 0.05
    desc_w     = 1.65
    chart_x    = desc_x + desc_w + 0.05
    chart_w    = 2.20
    bar_h      = 0.20
    bar_gap    = 0.04
    t_max      = max(t[2] for t in turns)

    for i, (turn, desc, ms, col) in enumerate(turns):
        y = chart_y0 + i * (bar_h + bar_gap)
        w = max(0.10, (ms / t_max) * chart_w)
        add_text(s, label_x, y, label_w, bar_h, turn,
                 size=10, bold=True, color=DARK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc_x, y, desc_w, bar_h, desc,
                 size=8, italic=True, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, chart_x, y, w, bar_h, fill=col)
        # value label
        if ms < 1:
            v_str = "< 1 ms"
        else:
            v_str = f"{ms:.1f} ms"
        add_text(s, chart_x + w + 0.05, y, 0.85, bar_h, v_str,
                 size=10, color=DARK, bold=(ms < 1),
                 anchor=MSO_ANCHOR.MIDDLE)

    # ── Right: takeaway callout (aligned with chart vertical span) ───────────
    add_rect(s, 7.10, 3.45, 2.40, 1.50, fill=NAVY)
    add_text(s, 7.10, 3.53, 2.40, 0.24,
             "EFFECT", size=10, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER)
    add_text(s, 7.10, 3.85, 2.40, 0.55,
             "24 ms → < 1 ms",
             size=22, bold=True, color=WHITE, font=HEADER_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.10, 4.45, 2.40, 0.50,
             "by turn 4, with the same\n5 000 iter/turn requested",
             size=9, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)

    # ── Bottom: two impact panels (raw overhead disclosure + strategic value) ─
    strip_y = 5.05
    strip_h = 0.45

    # Left — raw cost disclosure
    add_rect(s, 0.5, strip_y, 5.70, strip_h, fill=ICE)
    add_rect(s, 0.5, strip_y, 0.10, strip_h, fill=DEEP)
    add_text(s, 0.70, strip_y + 0.04, 5.40, 0.20,
             "Honest cost", size=9, bold=True, color=NAVY)
    add_text(s, 0.70, strip_y + 0.23, 5.40, 0.20,
             "Solver bookkeeping costs 1-25 % iter/s vs plain UCT; "
             "reuse adds an O(k) BFS compaction once per move.",
             size=8, color=SLATE, italic=True)

    # Right — strategic value
    add_rect(s, 6.30, strip_y, 3.20, strip_h, fill=ICE)
    add_rect(s, 6.30, strip_y, 0.10, strip_h, fill=FOREST)
    add_text(s, 6.50, strip_y + 0.04, 2.95, 0.20,
             "Strategic value", size=9, bold=True, color=NAVY)
    add_text(s, 6.50, strip_y + 0.23, 2.95, 0.20,
             "Provably correct endgame play  +  deeper effective search.",
             size=8, color=SLATE, italic=True)

    # ── Save ─────────────────────────────────────────────────────────────────
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"  1 slide · {OUT_PATH.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
