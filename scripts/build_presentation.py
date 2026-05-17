"""Build the PopOut AI defence presentation (docs/presentation.pptx).

Produces an 18-slide deck for the AI 2025/2026 oral defence (≤10 min budget).

Design system:
  - Palette "Ocean Gradient": navy / deep-blue / teal + accents (amber, coral,
    forest green) for emphasis.
  - Header font: Cambria (academic).  Body font: Calibri (clean).
  - Consistent header bar with section chip + footer with page numbers.
  - Each slide has at least one visual element (diagram, chart, KPI, figure).
  - Two cover slides (title + conclusion) on dark navy; content slides white.

Usage:
    python scripts/build_presentation.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parents[1]
FIG_DIR  = ROOT / "data/figures"
OUT_PATH = ROOT / "docs/presentation.pptx"

# ─── Palette ──────────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x21, 0x29, 0x5C)   # primary dark
NAVY_DK  = RGBColor(0x12, 0x1E, 0x3F)   # near-black accent
DEEP     = RGBColor(0x06, 0x5A, 0x82)   # primary mid
TEAL     = RGBColor(0x1C, 0x72, 0x93)   # accent
TEAL_LT  = RGBColor(0x5A, 0xA4, 0xB8)   # softer teal for fills
ICE      = RGBColor(0xE6, 0xF1, 0xF5)   # subtle bg fill
ICE_DK   = RGBColor(0xC4, 0xDB, 0xE3)   # slightly darker ice for variation
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x12, 0x1E, 0x3F)
SLATE    = RGBColor(0x33, 0x44, 0x55)   # body text
GRAY     = RGBColor(0x6B, 0x7A, 0x88)   # muted text
LIGHT_GRAY = RGBColor(0xB5, 0xC1, 0xCA)
AMBER    = RGBColor(0xF4, 0xA3, 0x41)   # warm accent for KPIs
CORAL    = RGBColor(0xE5, 0x71, 0x4E)   # warning accent
FOREST   = RGBColor(0x2E, 0x8B, 0x57)   # success accent
RED      = RGBColor(0xC1, 0x3F, 0x2F)

HEADER_FONT = "Cambria"
BODY_FONT   = "Calibri"
MONO_FONT   = "Consolas"

PAGE_W = 10.0
PAGE_H = 5.625


# ─── Primitives ──────────────────────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *,
             size=14, bold=False, italic=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *,
                size=13, color=SLATE, bullet_color=TEAL,
                bullet=True, spacing=6, font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0.0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        if bullet:
            b = p.add_run()
            b.text = "•  "
            b.font.name = font
            b.font.size = Pt(size)
            b.font.bold = True
            b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.5,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_circle(slide, x, y, d, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                 Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_arrow(slide, x1, y1, x2, y2, *, color=TEAL, width=2):
    """Simple horizontal/vertical arrow as a line+triangle approximation."""
    conn = slide.shapes.add_connector(2, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


# ─── Slide-level helpers ──────────────────────────────────────────────────────
def header(slide, title, *, chip=None, subtitle=None):
    """Standard header for content slides: chip + title + thin rule."""
    # Background top band (white slides have white, this is just for spacing)
    if chip:
        add_rect(slide, 0.5, 0.32, 0.65, 0.33, fill=TEAL)
        add_text(slide, 0.5, 0.32, 0.65, 0.33, chip,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=BODY_FONT)
        add_text(slide, 1.30, 0.28, 8.2, 0.45, title,
                 size=22, bold=True, color=NAVY, font=HEADER_FONT)
    else:
        add_text(slide, 0.5, 0.28, 9.0, 0.45, title,
                 size=22, bold=True, color=NAVY, font=HEADER_FONT)
    if subtitle:
        add_text(slide, 1.30 if chip else 0.5, 0.72, 8.2, 0.28, subtitle,
                 size=11, italic=True, color=GRAY)
    # subtle rule
    add_rect(slide, 0.5, 1.05, 9.0, 0.015, fill=ICE_DK)


def footer(slide, page, total=18):
    add_text(slide, 0.5, 5.34, 4.5, 0.22,
             "PopOut AI  ·  IA 2025/2026", size=9, color=LIGHT_GRAY)
    add_text(slide, 5.0, 5.34, 4.5, 0.22,
             f"{page} / {total}", size=9, color=LIGHT_GRAY,
             align=PP_ALIGN.RIGHT)


def kpi_card(slide, x, y, w, h, value, label, *, accent=TEAL,
             value_size=36, label_size=11):
    """Big-number card."""
    add_rect(slide, x, y, w, h, fill=ICE)
    add_rect(slide, x, y, w, 0.08, fill=accent)
    add_text(slide, x, y + 0.20, w, h * 0.45, value,
             size=value_size, bold=True, color=accent,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)
    add_text(slide, x + 0.1, y + h * 0.65, w - 0.2, h * 0.30, label,
             size=label_size, color=SLATE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, title, bullets, *, accent=TEAL,
         title_size=14, body_size=12):
    """Card with accent top stripe + title + bullets."""
    add_rect(slide, x, y, w, h, fill=WHITE, line=ICE_DK, line_w=0.5)
    add_rect(slide, x, y, w, 0.08, fill=accent)
    add_text(slide, x + 0.15, y + 0.18, w - 0.3, 0.35, title,
             size=title_size, bold=True, color=NAVY)
    add_bullets(slide, x + 0.15, y + 0.55, w - 0.3, h - 0.65,
                bullets, size=body_size, color=SLATE,
                bullet_color=accent, spacing=4)


def step_circle(slide, cx, cy, d, number, *, color=TEAL):
    """Numbered circle (used for process diagrams)."""
    x = cx - d / 2
    y = cy - d / 2
    add_circle(slide, x, y, d, fill=color)
    add_text(slide, x, y, d, d, str(number),
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)


# ─── Build ────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    blank = prs.slide_layouts[6]

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — COVER
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    # Left teal stripe
    add_rect(s, 0.0, 0.0, 0.30, PAGE_H, fill=TEAL)

    # Decorative discs (Connect-4 reference) in bottom-right corner
    disc_y = 4.30
    colors_disc = [TEAL, AMBER, WHITE, AMBER, TEAL, WHITE, AMBER]
    for i in range(7):
        add_circle(s, 8.10 + (i % 4) * 0.36, disc_y + (i // 4) * 0.36,
                   0.30, fill=colors_disc[i % len(colors_disc)])

    # Course tag
    add_text(s, 0.85, 0.95, 8.5, 0.32,
             "INTELIGÊNCIA ARTIFICIAL  ·  2025 / 2026",
             size=13, color=TEAL_LT, bold=True, font=BODY_FONT)

    # Title (big)
    add_text(s, 0.85, 1.40, 8.5, 1.05,
             "PopOut AI",
             size=64, bold=True, color=WHITE, font=HEADER_FONT)

    # Subtitle
    add_text(s, 0.85, 2.55, 8.5, 0.50,
             "MCTS + ID3 para uma variante do Connect-4",
             size=20, italic=True, color=ICE, font=BODY_FONT)

    # Accent rule
    add_rect(s, 0.85, 3.18, 1.6, 0.03, fill=AMBER)

    # Group
    add_text(s, 0.85, 3.35, 8.0, 0.30,
             "GRUPO", size=11, bold=True, color=TEAL_LT)
    add_text(s, 0.85, 3.65, 8.0, 0.30,
             "Duarte Meneses dos Santos Sousa Gomes   ·   202409386",
             size=12, color=WHITE)
    add_text(s, 0.85, 3.95, 8.0, 0.30,
             "José Paulo Pacheco de Sousa   ·   202405046",
             size=12, color=WHITE)
    add_text(s, 0.85, 4.25, 8.0, 0.30,
             "Tiago Braga da Cruz Frada de Sousa   ·   202405406",
             size=12, color=WHITE)

    add_text(s, 0.85, 5.10, 8.0, 0.30,
             "github.com/Difl4/popout-ai   ·   Maio 2026",
             size=10, color=TEAL_LT, italic=True)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — AGENDA
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Agenda", chip="00")

    sections = [
        ("01", "Problema & Requisitos",     "PopOut + 3 cenários",       TEAL),
        ("02", "Arquitectura",              "Engine · Agentes · UI",     DEEP),
        ("03", "Pesquisa Adversarial",      "MCTS · UCT · Solver",       NAVY),
        ("04", "Decision Trees",            "ID3 · Iris · PopOut",       AMBER),
        ("05", "Resultados & Dificuldades", "Métricas · Trade-offs",     CORAL),
        ("06", "Conclusão",                 "Síntese · Trabalho futuro", FOREST),
    ]
    cols = 3
    cw, ch = 2.95, 1.85
    for i, (num, title, sub, col) in enumerate(sections):
        x = 0.5 + (i % cols) * (cw + 0.08)
        y = 1.30 + (i // cols) * (ch + 0.15)
        add_rect(s, x, y, cw, ch, fill=ICE)
        add_rect(s, x, y, 0.10, ch, fill=col)
        add_text(s, x + 0.25, y + 0.20, cw - 0.35, 0.40, num,
                 size=22, bold=True, color=col, font=HEADER_FONT)
        add_text(s, x + 0.25, y + 0.70, cw - 0.35, 0.45, title,
                 size=14, bold=True, color=NAVY, font=HEADER_FONT)
        add_text(s, x + 0.25, y + 1.15, cw - 0.35, 0.55, sub,
                 size=11, color=GRAY, italic=True)

    footer(s, 2)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — O PROBLEMA: POPOUT
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "O Problema: PopOut",
           chip="01", subtitle="Connect-4 com uma jogada adicional")

    # Left: rules
    add_text(s, 0.5, 1.30, 5.0, 0.40,
             "Mecânica do jogo", size=15, bold=True, color=DEEP)
    add_bullets(s, 0.5, 1.75, 5.0, 1.5, [
        "Tabuleiro 7 × 6, dois jogadores alternam.",
        "DROP: largar peça por cima da coluna.",
        "POP: remover peça própria do fundo — peças acima descem.",
        "Vencedor: 4 alinhadas (horizontal / vertical / diagonal).",
    ], size=13)

    add_text(s, 0.5, 3.30, 5.0, 0.40,
             "Três regras especiais", size=15, bold=True, color=DEEP)
    add_bullets(s, 0.5, 3.75, 5.0, 1.4, [
        "Pop simultâneo (4-em-linha para ambos) → ganha quem fez o pop.",
        "Tabuleiro cheio: o jogador a mover pode declarar empate.",
        "Repetição tripla: qualquer um pode declarar empate.",
    ], size=12, color=SLATE)

    # Right: stylized board mini-diagram
    bx, by = 6.20, 1.30
    board_w, board_h = 3.20, 2.75
    add_rect(s, bx, by, board_w, board_h, fill=NAVY)
    # 7 cols × 6 rows of empty discs
    cell_size = 0.36
    pad_x = (board_w - 7 * cell_size) / 2
    pad_y = (board_h - 6 * cell_size) / 2
    # Set up a small position for visual
    pos = [
        # (row, col, player) — sample non-trivial position
        (5, 1, 1), (5, 2, 1), (5, 3, 1), (5, 5, 2),
        (4, 2, 2), (4, 3, 1),
        (3, 3, 1),
    ]
    pos_map = {(r, c): p for (r, c, p) in pos}
    for r in range(6):
        for c in range(7):
            cx = bx + pad_x + c * cell_size
            cy = by + pad_y + r * cell_size
            who = pos_map.get((r, c))
            if who == 1:
                color = AMBER
            elif who == 2:
                color = WHITE
            else:
                color = NAVY_DK
            add_circle(s, cx + 0.04, cy + 0.04, cell_size - 0.08, fill=color)

    # POP arrow indicating the POP move
    add_text(s, bx, by + board_h + 0.10, board_w, 0.30,
             "↑  POP", size=14, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER, font=HEADER_FONT)

    # Board caption
    add_text(s, bx, by + board_h + 0.45, board_w, 0.50,
             "DROP em qualquer coluna, ou POP\nremovendo peça própria do fundo.",
             size=10, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    footer(s, 3)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — REQUISITOS DO GUIÃO
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Requisitos do Guião",
           chip="01", subtitle="§4 do IA_2526_Trab.pdf")

    # Left column — game scenarios
    add_text(s, 0.5, 1.30, 4.4, 0.40,
             "Cenários de jogo (§3)", size=15, bold=True, color=DEEP)
    scenarios = [
        ("Humano vs Humano",         "GUI + CLI"),
        ("Humano vs Computador",     "GUI: 13 dificuldades  ·  CLI"),
        ("Computador vs Computador", "GUI Arena  ·  CLI torneio"),
    ]
    for i, (name, where) in enumerate(scenarios):
        y = 1.85 + i * 0.70
        add_circle(s, 0.55, y, 0.30, fill=FOREST)
        add_text(s, 0.55, y, 0.30, 0.30, "✓",
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 1.00, y - 0.02, 3.6, 0.35, name,
                 size=13, bold=True, color=DARK)
        add_text(s, 1.00, y + 0.28, 3.6, 0.30, where,
                 size=10, color=GRAY, italic=True)

    # Right column — algorithmic
    add_text(s, 5.20, 1.30, 4.3, 0.40,
             "Componentes algorítmicos (§4.1 / §4.2)",
             size=15, bold=True, color=DEEP)
    algos = [
        "MCTS com UCT",
        "Diferentes números de filhos (TopK-UCT)",
        "Variantes além do standard (Solver, Numba)",
        "ID3 from scratch  ·  sem scikit-learn",
        "Iris (com discretização)",
        "Dataset PopOut a partir do MCTS",
    ]
    for i, item in enumerate(algos):
        y = 1.85 + i * 0.40
        add_circle(s, 5.25, y, 0.22, fill=FOREST)
        add_text(s, 5.25, y, 0.22, 0.22, "✓",
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 5.60, y - 0.02, 4.0, 0.32, item,
                 size=12, color=DARK)

    # Bottom callout
    add_rect(s, 0.5, 4.85, 9.0, 0.40, fill=NAVY)
    add_text(s, 0.5, 4.85, 9.0, 0.40,
             "Todos os requisitos cumpridos  ·  4 critérios: 30% adversarial  /  30% trees  /  30% técnica  /  10% comunicação",
             size=11, color=ICE, italic=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 4)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — ARQUITECTURA + CONTRIBUIÇÕES
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Arquitectura & Contribuições", chip="02")

    # Three-layer architecture stack
    layer_x, layer_w = 0.5, 5.4
    layer_h = 0.90
    layers = [
        ("Interface",     "CLI  ·  GUI Pygame  ·  GUI Arena",                NAVY),
        ("Agentes",       "MCTS variantes  ·  ID3 from scratch  ·  Híbridos", DEEP),
        ("Engine",        "Bitboard 7×6  ·  Regras (3 especiais)  ·  Numba kernels", TEAL),
    ]
    for i, (name, body, col) in enumerate(layers):
        y = 1.30 + i * (layer_h + 0.10)
        add_rect(s, layer_x, y, layer_w, layer_h, fill=ICE)
        add_rect(s, layer_x, y, 0.12, layer_h, fill=col)
        add_text(s, layer_x + 0.25, y + 0.13, layer_w - 0.35, 0.32, name,
                 size=14, bold=True, color=col, font=HEADER_FONT)
        add_text(s, layer_x + 0.25, y + 0.45, layer_w - 0.35, 0.40, body,
                 size=11, color=SLATE)

    # Right column — contributions
    cx, cy = 6.20, 1.30
    add_text(s, cx, cy, 3.30, 0.35,
             "Distribuição de tarefas", size=14, bold=True, color=NAVY,
             font=HEADER_FONT)

    contribs = [
        ("Duarte",  "Bitboard  ·  MCTS-Solver  ·\nKernels Numba  ·  Pipeline ID3", TEAL),
        ("José",    "ID3 from scratch  ·  Dataset\ngenerator  ·  Análise leakage", AMBER),
        ("Tiago",   "GUI Pygame  ·  CLI/torneio  ·\nBenchmarks  ·  Testes",        DEEP),
    ]
    for i, (name, role, col) in enumerate(contribs):
        y = cy + 0.50 + i * 1.05
        add_rect(s, cx, y, 3.30, 0.95, fill=ICE)
        add_rect(s, cx, y, 0.08, 0.95, fill=col)
        add_text(s, cx + 0.20, y + 0.10, 3.0, 0.32, name,
                 size=14, bold=True, color=col, font=HEADER_FONT)
        add_text(s, cx + 0.20, y + 0.42, 3.0, 0.55, role,
                 size=10, color=SLATE)

    footer(s, 5)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — MCTS: 4 FASES
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "MCTS: As 4 Fases Canónicas",
           chip="03", subtitle="Cada iteração percorre selection → expansion → simulation → backprop")

    phases = [
        ("Selection",
         "UCT escolhe o ramo\ncom maior promessa.",
         TEAL),
        ("Expansion",
         "Adiciona um novo nó\ndepois da folha actual.",
         DEEP),
        ("Simulation",
         "Rollout aleatório até\nfim ou profundidade limite.",
         NAVY),
        ("Backpropagation",
         "Actualiza visitas e\nvalor ao subir.",
         AMBER),
    ]
    cw = 2.20
    cy = 1.35
    ch = 3.55
    gap = 0.10
    for i, (name, desc, col) in enumerate(phases):
        x = 0.5 + i * (cw + gap)
        add_rect(s, x, cy, cw, ch, fill=ICE)
        add_rect(s, x, cy, cw, 0.08, fill=col)
        # Big step number
        step_circle(s, x + cw / 2, cy + 0.85, 0.85, i + 1, color=col)
        # Name
        add_text(s, x + 0.10, cy + 1.50, cw - 0.20, 0.40, name,
                 size=14, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font=HEADER_FONT)
        # Description
        add_text(s, x + 0.15, cy + 1.95, cw - 0.30, 1.20, desc,
                 size=11, color=SLATE, align=PP_ALIGN.CENTER)

    # Arrow underline indicating circular nature
    add_text(s, 0.5, 5.00, 9.0, 0.25,
             "→ repete milhares a milhões de vezes →",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    footer(s, 6)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — UCT + VARIANTES
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "UCT + Variantes",
           chip="03", subtitle="Exploração vs exploitation com prova de convergência")

    # Left — UCT formula panel
    add_rect(s, 0.5, 1.30, 4.20, 3.55, fill=NAVY)
    add_text(s, 0.5, 1.45, 4.20, 0.40,
             "Fórmula UCT", size=13, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER, font=HEADER_FONT)
    add_text(s, 0.5, 1.95, 4.20, 0.85,
             "Q  +  C · √(ln N / n)",
             size=26, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)
    add_rect(s, 1.5, 2.95, 2.2, 0.02, fill=TEAL)

    add_text(s, 0.7, 3.10, 3.8, 0.30,
             "Q  ·  taxa de vitória média do filho",
             size=11, color=ICE)
    add_text(s, 0.7, 3.45, 3.8, 0.30,
             "N  ·  visitas do nó pai",
             size=11, color=ICE)
    add_text(s, 0.7, 3.80, 3.8, 0.30,
             "n  ·  visitas do filho",
             size=11, color=ICE)
    add_text(s, 0.7, 4.15, 3.8, 0.30,
             "C  ·  √2 ≈ 1.414  (Kocsis & Szepesvári)",
             size=11, color=ICE, italic=True)
    add_text(s, 0.5, 4.55, 4.20, 0.25,
             "Visit unvisited first  →  score = +∞",
             size=10, italic=True, color=AMBER,
             align=PP_ALIGN.CENTER)

    # Right — variants table
    add_text(s, 5.0, 1.30, 4.5, 0.35,
             "Cinco variantes implementadas",
             size=14, bold=True, color=NAVY, font=HEADER_FONT)

    variants = [
        ("StandardUCT",       "Baseline canónico",                    TEAL),
        ("ExperimentalUCT",   "Prioriza não-visitados",               DEEP),
        ("TopK-UCT",          "Limita a k filhos  (§4.1)",            AMBER),
        ("MCTS-Solver",       "Propaga provas matemáticas",           CORAL),
        ("Numba (4 var.)",    "Kernels @njit  ·  ~220 k iter/s",      FOREST),
    ]
    for i, (name, desc, col) in enumerate(variants):
        y = 1.80 + i * 0.58
        add_rect(s, 5.0, y, 4.5, 0.48, fill=ICE)
        add_rect(s, 5.0, y, 0.10, 0.48, fill=col)
        add_text(s, 5.25, y + 0.05, 1.85, 0.42, name,
                 size=12, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 7.10, y + 0.05, 2.35, 0.42, desc,
                 size=11, color=SLATE,
                 anchor=MSO_ANCHOR.MIDDLE)

    footer(s, 7)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — K-CHILDREN ANALYSIS
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Análise do Número de Filhos (k)",
           chip="03", subtitle="Requisito §4.1 — explorar diferentes números de filhos seleccionados")

    # Left column (filled, balanced height with right figure)
    add_text(s, 0.5, 1.30, 4.5, 0.40,
             "Variante TopK-UCT", size=14, bold=True, color=DEEP)
    add_bullets(s, 0.5, 1.75, 4.5, 1.30, [
        "Após cada filho visitado pelo menos 1×, restringe a "
        "selecção aos K filhos com maior UCT.",
        "Reduz ramos com baixa Q sem cortar exploração inicial.",
        "Hiperparâmetro k controla a largura da pesquisa.",
    ], size=12)

    add_text(s, 0.5, 3.15, 4.5, 0.40,
             "Metodologia", size=14, bold=True, color=DEEP)
    add_bullets(s, 0.5, 3.60, 4.5, 1.10, [
        "15 estados aleatórios  ·  600 iter por agente",
        "k ∈ {3, 5, 7, 14}  (todos)",
        "Métrica: consistência da jogada óptima",
    ], size=11, color=SLATE)

    # Highlighted finding box
    add_rect(s, 0.5, 4.85, 4.5, 0.35, fill=AMBER)
    add_text(s, 0.5, 4.85, 4.5, 0.35,
             "k = 7 já preserva 100 % da consistência do baseline.",
             size=11, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Right — figure (aspect ratio 2.9:1, fits within content width)
    fig = FIG_DIR / "mcts_k_children.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(5.20), Inches(1.30),
                             width=Inches(4.30))
    else:
        add_rect(s, 5.20, 1.30, 4.30, 1.50, fill=ICE)
        add_text(s, 5.20, 1.95, 4.30, 0.30, "(mcts_k_children.png)",
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Right — insights below figure (balances the left column visually)
    add_text(s, 5.20, 3.00, 4.30, 0.35,
             "Resultados observados", size=14, bold=True, color=DEEP)
    add_bullets(s, 5.20, 3.40, 4.30, 1.40, [
        "Consistência estável em k ∈ {3, 5, 7, 10, 14}.",
        "Concordância com k=14 (standard) próxima de 33 %.",
        "Custo computacional cresce linearmente com k.",
    ], size=11, color=SLATE, spacing=3)

    footer(s, 8)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — MCTS-SOLVER
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "MCTS-Solver: do Probabilístico ao Demonstrado",
           chip="03",
           subtitle="Propagação AND/OR de WIN / LOSS / DRAW + distância minimax")

    # Left: rules of propagation
    add_text(s, 0.5, 1.25, 4.7, 0.40,
             "Regras de propagação", size=14, bold=True, color=DEEP)

    rules = [
        ("WIN",   "∃ filho LOSS",         FOREST),
        ("LOSS",  "TODOS filhos WIN  ·  fully expanded",  RED),
        ("DRAW",  "todos DRAW ou consenso",                AMBER),
    ]
    for i, (label, cond, col) in enumerate(rules):
        y = 1.75 + i * 0.65
        add_rect(s, 0.5, y, 1.20, 0.50, fill=col)
        add_text(s, 0.5, y, 1.20, 0.50, label,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=HEADER_FONT)
        add_text(s, 1.85, y, 3.30, 0.50, cond,
                 size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    # Distance minimax box
    add_rect(s, 0.5, 3.80, 4.7, 1.20, fill=ICE)
    add_text(s, 0.5, 3.85, 4.7, 0.32,
             "Distância minimax", size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, font=HEADER_FONT)
    add_text(s, 0.7, 4.20, 4.3, 0.30,
             "WIN  →  min( dist(filho) ) + 1   (win fast)",
             size=11, color=SLATE, font=MONO_FONT)
    add_text(s, 0.7, 4.55, 4.3, 0.30,
             "LOSS  →  max( dist(filho) ) + 1   (lose slow)",
             size=11, color=SLATE, font=MONO_FONT)

    # Right: mini-diagram showing tree with labels
    rx, ry = 5.50, 1.25
    # Root
    add_circle(s, rx + 1.65, ry, 0.45, fill=NAVY)
    add_text(s, rx + 1.65, ry, 0.45, 0.45, "?",
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + 1.65 + 0.50, ry + 0.10, 1.5, 0.25,
             "root  ·  to move", size=9, color=GRAY, italic=True)

    # Children level 1
    child_y = ry + 1.10
    children1 = [(rx + 0.10, "?"), (rx + 1.65, "L"), (rx + 3.20, "?")]
    for cx, label in children1:
        col = RED if label == "L" else NAVY
        add_circle(s, cx, child_y, 0.45, fill=col)
        add_text(s, cx, child_y, 0.45, 0.45, label,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_arrow(s, rx + 1.65 + 0.225, ry + 0.45,
                  cx + 0.225, child_y, color=LIGHT_GRAY)

    # Children level 2 (under the leftmost child)
    leaf_y = child_y + 1.10
    leaves = [(rx - 0.20, "W"), (rx + 0.40, "W")]
    for cx, label in leaves:
        col = FOREST if label == "W" else RED
        add_circle(s, cx, leaf_y, 0.40, fill=col)
        add_text(s, cx, leaf_y, 0.40, 0.40, label,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_arrow(s, rx + 0.325, child_y + 0.45,
                  cx + 0.20, leaf_y, color=LIGHT_GRAY)

    # Annotation
    add_text(s, 5.50, 3.65, 4.0, 0.30,
             "Filho LOSS → root provado WIN",
             size=11, italic=True, color=FOREST, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, 5.50, 4.00, 4.0, 0.30,
             "Folhas confirmadas propagam para cima.",
             size=10, color=GRAY, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(s, 5.50, 4.40, 4.0, 0.60,
             "Pesquisa pode terminar mais cedo:\nbasta uma jogada provada vencedora.",
             size=10, color=SLATE, align=PP_ALIGN.CENTER, italic=True)

    footer(s, 9)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — TREE REUSE + NUMBA
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Tree Reuse + Numba: Performance",
           chip="03", subtitle="Duas optimizações ortogonais — ~44× speedup do kernel + reuse entre turnos")

    # Top-left — Tree Reuse explanation
    add_text(s, 0.5, 1.25, 4.5, 0.35,
             "Tree Reuse  ·  reaproveita sub-árvore",
             size=13, bold=True, color=DEEP)

    # Before/After diagram (compaction)
    # Before
    add_text(s, 0.5, 1.65, 2.2, 0.25,
             "ANTES", size=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, 0.7, 1.95, 1.8, 1.15, fill=ICE)
    # Draw tree skeleton
    add_circle(s, 1.50, 2.05, 0.20, fill=NAVY)
    for i, dx in enumerate([-0.40, 0.0, 0.40]):
        add_circle(s, 1.50 + dx, 2.45, 0.18, fill=TEAL_LT)
    # Highlight the chosen subtree
    add_circle(s, 1.50, 2.45, 0.18, fill=AMBER)
    for i, dx in enumerate([-0.20, 0.20]):
        add_circle(s, 1.50 + dx, 2.85, 0.14, fill=AMBER)

    # Arrow
    add_arrow(s, 2.65, 2.55, 3.10, 2.55, color=NAVY_DK, width=3)

    # After
    add_text(s, 3.20, 1.65, 1.7, 0.25,
             "DEPOIS", size=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    add_rect(s, 3.30, 1.95, 1.55, 1.15, fill=ICE)
    add_circle(s, 4.07, 2.10, 0.20, fill=AMBER)
    for i, dx in enumerate([-0.20, 0.20]):
        add_circle(s, 4.07 + dx, 2.50, 0.16, fill=AMBER)

    add_text(s, 0.5, 3.20, 4.5, 0.40,
             "BFS compaction re-indexa nós para slots 0..k antes do kernel.",
             size=10, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Bottom: Numba performance chart (built with shapes)
    add_text(s, 5.20, 1.25, 4.3, 0.35,
             "Numba JIT  ·  throughput",
             size=13, bold=True, color=DEEP)

    # Bar chart manually (using rectangles)
    # Labels live to the LEFT of bars (always readable); values to the RIGHT.
    # Numbers measured on the project's own benchmark (Technical_Documentation
    # notebook). Reuse does NOT increase raw iter/s — its benefit is described
    # in the slide narrative below.
    bars = [
        ("StandardUCT",         5_000,    TEAL_LT),
        ("NumbaMCTS",           30_000,   TEAL),
        ("FlatNumbaMCTS",       220_000,  DEEP),
        ("FlatNumbaSolverMCTS", 220_000,  NAVY),
        ("ReuseFlat Solver",    220_000,  AMBER),
    ]
    import math
    bar_max = max(b[1] for b in bars)
    bar_y0  = 1.70
    bar_h   = 0.32
    bar_gap = 0.12
    label_x = 5.20
    label_w = 1.50
    chart_x = label_x + label_w + 0.05
    chart_w = 2.10
    for i, (name, val, col) in enumerate(bars):
        y = bar_y0 + i * (bar_h + bar_gap)
        w = max(0.10, math.log10(val + 1) / math.log10(bar_max + 1) * chart_w)
        # Left-side label (always readable)
        add_text(s, label_x, y, label_w, bar_h, name,
                 size=10, bold=True, color=DARK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # Bar
        add_rect(s, chart_x, y, w, bar_h, fill=col)
        # Right-side value
        if val >= 1_000_000:
            label_val = f"{val // 1_000_000} M"
        elif val >= 1000:
            label_val = f"{val // 1000} k"
        else:
            label_val = f"{val}"
        add_text(s, chart_x + w + 0.05, y, 0.95, bar_h,
                 f"{label_val} iter/s",
                 size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, 5.20, 4.50, 4.3, 0.25,
             "log-scale  ·  unidades: iter/s",
             size=9, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

    # Massive speedup callout
    add_rect(s, 0.5, 4.85, 4.5, 0.40, fill=AMBER)
    add_text(s, 0.5, 4.85, 4.5, 0.40,
             "≈ 44× speedup  vs  Python puro",
             size=13, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)

    footer(s, 10)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 11 — ID3 ALGORITMO
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "ID3: Decision Tree do Zero",
           chip="04",
           subtitle="Sem scikit-learn  ·  apenas numpy + pandas para data wrangling")

    # Three-step process
    steps = [
        ("01", "Entropia",
         "H = − Σ pᵢ · log₂(pᵢ)",
         "Mede a impureza\ndo conjunto.", TEAL),
        ("02", "Information Gain",
         "IG = H(S) − Σ |Sv|/|S| · H(Sv)",
         "Ganho de informação\nao dividir por feature.", DEEP),
        ("03", "Split Recursivo",
         "feature* = argmax IG",
         "Repete até nó puro,\nmax_depth, ou sem features.", AMBER),
    ]

    cw = 2.95
    cy = 1.25
    ch = 3.40
    for i, (num, name, formula, desc, col) in enumerate(steps):
        x = 0.5 + i * (cw + 0.08)
        add_rect(s, x, cy, cw, ch, fill=ICE)
        add_rect(s, x, cy, cw, 0.10, fill=col)
        # step number
        add_text(s, x + 0.15, cy + 0.20, 0.7, 0.45, num,
                 size=22, bold=True, color=col, font=HEADER_FONT)
        # name
        add_text(s, x + 0.95, cy + 0.25, cw - 1.0, 0.40, name,
                 size=15, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        # formula box
        add_rect(s, x + 0.15, cy + 0.85, cw - 0.30, 0.60, fill=WHITE,
                 line=col, line_w=0.75)
        add_text(s, x + 0.15, cy + 0.85, cw - 0.30, 0.60, formula,
                 size=12, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=MONO_FONT)
        # description
        add_text(s, x + 0.20, cy + 1.65, cw - 0.40, 1.30, desc,
                 size=11, color=SLATE, align=PP_ALIGN.CENTER)

    # Bottom callout
    add_text(s, 0.5, 4.85, 9.0, 0.30,
             "src/decision_tree/id3/learner.py  ·  100% próprio  ·  47 testes a passar",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    footer(s, 11)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 12 — IRIS WARM-UP
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Iris (Warm-up)",
           chip="04", subtitle="Validar o ID3 antes do dataset PopOut")

    # Left — methodology
    add_text(s, 0.5, 1.30, 4.5, 0.35,
             "Metodologia", size=13, bold=True, color=DEEP)
    add_bullets(s, 0.5, 1.65, 4.5, 1.50, [
        "Discretização quantile  (q = 3 bins por feature)",
        "5 × 10-fold cross-validation  →  50 runs independentes",
        "max_depth = 10  ·  sem post-pruning",
        "Discretizer fit por fold (sem leakage)",
    ], size=11)

    # KPIs
    kpi_card(s, 0.5, 3.30, 2.10, 1.55,
             "95.07 %", "média 50 folds",
             accent=FOREST, value_size=24, label_size=10)
    kpi_card(s, 2.75, 3.30, 2.25, 1.55,
             "± 5.63 %", "desvio padrão",
             accent=DEEP, value_size=24, label_size=10)

    # Right — tree visualization
    fig = FIG_DIR / "iris_id3_tree.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(5.30), Inches(1.30),
                             width=Inches(4.20))
        add_text(s, 5.30, 4.75, 4.20, 0.30,
                 "Árvore ID3 — exemplo de 1 fold",
                 size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
    else:
        add_rect(s, 5.30, 1.30, 4.20, 3.55, fill=ICE)
        add_text(s, 5.30, 2.95, 4.20, 0.30, "(iris_id3_tree.png)",
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)

    footer(s, 12)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 13 — DATASET POPOUT
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Geração do Dataset PopOut",
           chip="04",
           subtitle="MCTS-Solver de 100k iterações etiqueta cada posição")

    # Top: pipeline diagram (left → right flow)
    flow_y = 1.30
    flow_h = 0.85
    boxes = [
        ("Self-play\n+ blunder tiers",      TEAL),
        ("Oracle\nReuseSolver 100k",        DEEP),
        ("Mirror\nhorizontal",              AMBER),
        ("CSV\n193 k posições",             FOREST),
    ]
    bw = 2.10
    bx0 = 0.5
    gap = 0.20
    for i, (text, col) in enumerate(boxes):
        x = bx0 + i * (bw + gap)
        add_rect(s, x, flow_y, bw, flow_h, fill=col)
        add_text(s, x, flow_y, bw, flow_h, text,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 font=HEADER_FONT)
        # Chevron arrow between boxes — much clearer than a line connector
        if i < len(boxes) - 1:
            ar_x = x + bw + 0.01
            ar_y = flow_y + flow_h * 0.25
            add_rect(s, ar_x, ar_y, gap - 0.02, flow_h * 0.50,
                     shape=MSO_SHAPE.RIGHT_ARROW, fill=NAVY_DK)

    # KPIs row
    kpi_w = 2.10
    kpi_y = 2.55
    kpi_h = 1.30
    kpis = [
        ("193 k",    "posições",        TEAL),
        ("16 cores", "geração paralela", DEEP),
        ("~ 5 min",  "tempo total",     AMBER),
        ("39 %",     "provadas",        FOREST),
    ]
    for i, (val, lab, col) in enumerate(kpis):
        x = 0.5 + i * (kpi_w + 0.20)
        kpi_card(s, x, kpi_y, kpi_w, kpi_h, val, lab,
                 accent=col, value_size=22, label_size=10)

    # Bottom note
    add_text(s, 0.5, 4.05, 9.0, 0.35,
             "Diversidade",
             size=13, bold=True, color=DEEP, align=PP_ALIGN.CENTER)
    add_text(s, 0.5, 4.40, 9.0, 0.55,
             "Aberturas forçadas (7 × 7 = 49 combos)   ·   3 tiers de blunder (5 % / 25 % / 50 %)   ·   Mirroring duplica gratuitamente.",
             size=11, color=SLATE, italic=True, align=PP_ALIGN.CENTER)

    footer(s, 13)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 14 — DATA LEAKAGE (CRITICAL SLIDE)
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Rigor Metodológico: Detecção do Data Leakage",
           chip="04",
           subtitle="Auditoria do pipeline expôs duas camadas de leakage train/test")

    # Comparison table
    headers = ["Split", "State overlap", "Test acc", "Diagnóstico"]
    rows = [
        ("BUGGY (notebook original)",  "94.9 %", "0.876", "Bug do reset_index",         CORAL),
        ("RANDOM (row-level)",         "80.4 %", "0.752", "Duplicação intrínseca",      AMBER),
        ("GROUPED (state)",            "0.0 %",  "0.441", "Generalização verdadeira",   FOREST),
    ]
    col_x  = [0.5, 3.40, 5.40, 6.95]
    col_w  = [2.85, 2.0, 1.50, 2.55]

    # Header row
    hdr_y = 1.30
    add_rect(s, 0.5, hdr_y, 9.0, 0.42, fill=NAVY)
    for i, h in enumerate(headers):
        add_text(s, col_x[i] + 0.05, hdr_y, col_w[i] - 0.1, 0.42, h,
                 size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Data rows
    for r_i, (a, b, c, d, col) in enumerate(rows):
        y = hdr_y + 0.50 + r_i * 0.55
        fill = ICE if r_i % 2 == 0 else WHITE
        add_rect(s, 0.5, y, 9.0, 0.52, fill=fill)
        add_rect(s, 0.5, y, 0.10, 0.52, fill=col)
        add_text(s, col_x[0] + 0.05, y, col_w[0] - 0.1, 0.52, a,
                 size=12, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, col_x[1] + 0.05, y, col_w[1] - 0.1, 0.52, b,
                 size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE,
                 font=MONO_FONT)
        add_text(s, col_x[2] + 0.05, y, col_w[2] - 0.1, 0.52, c,
                 size=15, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE, font=MONO_FONT)
        add_text(s, col_x[3] + 0.05, y, col_w[3] - 0.1, 0.52, d,
                 size=11, color=GRAY, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Big conclusion bar
    add_rect(s, 0.5, 3.45, 9.0, 0.55, fill=FOREST)
    add_text(s, 0.5, 3.45, 9.0, 0.55,
             "Generalização real:  44 %   (vs 88 % memorização)",
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)

    # Baselines reference + honest framing
    add_text(s, 0.5, 4.20, 9.0, 0.30,
             "Baselines  ·  random: 7 %    ·   always-`drop_3`: 28 %    ·   ID3 (GROUPED): 44 %",
             size=11, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, 0.5, 4.55, 9.0, 0.45,
             "Reportamos os três splits no notebook — honestidade técnica > apresentação inflada.",
             size=11, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

    footer(s, 14)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 15 — AGENTE HÍBRIDO
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Agente Híbrido: Árvore + Tactical Safeguards",
           chip="04",
           subtitle="A árvore sozinha falha em táctica — o híbrido é a configuração real de jogo")

    # Left — flow diagram
    fy0 = 1.30
    fh = 0.55
    fgap = 0.20
    steps_flow = [
        ("can_win?",                "Joga vitória imediata.",       FOREST),
        ("opp_wins_next?",          "Bloqueia ameaça do oponente.", CORAL),
        ("else  →  tree.predict()", "Consulta a árvore ID3.",       DEEP),
    ]
    for i, (cond, action, col) in enumerate(steps_flow):
        y = fy0 + i * (fh + fgap)
        add_rect(s, 0.5, y, 4.5, fh, fill=ICE)
        add_rect(s, 0.5, y, 0.12, fh, fill=col)
        add_text(s, 0.70, y + 0.05, 1.5, fh - 0.10, cond,
                 size=12, bold=True, color=col, font=MONO_FONT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, 2.35, y + 0.05, 2.6, fh - 0.10, action,
                 size=11, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps_flow) - 1:
            add_arrow(s, 2.75, y + fh + 0.02,
                      2.75, y + fh + fgap - 0.02,
                      color=NAVY_DK, width=2)

    # Right — why this matters
    add_text(s, 5.50, 1.30, 4.0, 0.40,
             "Porquê?", size=14, bold=True, color=DEEP)

    add_rect(s, 5.50, 1.75, 4.0, 1.50, fill=ICE)
    add_text(s, 5.50, 1.85, 4.0, 0.35,
             "Tactical accuracy da árvore PURA",
             size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, 5.50, 2.20, 1.95, 0.95, "27 %",
             size=30, bold=True, color=CORAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)
    add_text(s, 5.50, 3.10, 1.95, 0.20, "opp_wins_next=1",
             size=9, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, 7.50, 2.20, 1.95, 0.95, "14 %",
             size=30, bold=True, color=CORAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font=HEADER_FONT)
    add_text(s, 7.50, 3.10, 1.95, 0.20, "pop moves",
             size=9, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Honesty statement
    add_rect(s, 5.50, 3.50, 4.0, 1.40, fill=AMBER)
    add_text(s, 5.65, 3.60, 3.7, 0.32,
             "Honestidade académica", size=12, bold=True, color=NAVY)
    add_text(s, 5.65, 3.95, 3.7, 0.90,
             "Os win-rates reportados são do HÍBRIDO (tree + lookahead), nunca da árvore isolada. Isto está declarado na cell 0 do notebook.",
             size=10, color=DARK, italic=True)

    footer(s, 15)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 16 — RESULTADOS GLOBAIS
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Resultados Globais",
           chip="05",
           subtitle="ID3 híbrido em paridade com Solver 100k a uma fracção do tempo")

    # KPIs row
    kpi_y = 1.25
    kpi_h = 1.45
    kpis = [
        ("220 k",   "iter/s\n(FlatNumba)",        TEAL),
        ("44 %",    "generalização\n(ID3 GROUPED)", DEEP),
        ("52 %",    "win-rate vs\nSolver 100k",   AMBER),
        ("< 2 ms",  "inference ID3\npor jogada",  FOREST),
    ]
    kw = 2.15
    for i, (val, lab, col) in enumerate(kpis):
        x = 0.5 + i * (kw + 0.10)
        kpi_card(s, x, kpi_y, kw, kpi_h, val, lab,
                 accent=col, value_size=28, label_size=10)

    # Figure — constrain by HEIGHT (available: 5.30 − 2.95 = 2.35")
    fig = FIG_DIR / "speed_quality_tradeoff.png"
    if fig.exists():
        s.shapes.add_picture(str(fig), Inches(0.5), Inches(2.95),
                             height=Inches(2.30))
    else:
        add_rect(s, 0.5, 2.95, 5.6, 2.30, fill=ICE)
        add_text(s, 0.5, 4.00, 5.6, 0.30,
                 "(speed_quality_tradeoff.png)",
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # Right notes
    add_text(s, 6.40, 2.95, 3.1, 0.35,
             "Conclusões empíricas", size=13, bold=True, color=NAVY,
             font=HEADER_FONT)
    add_bullets(s, 6.40, 3.35, 3.1, 1.85, [
        "ID3 híbrido paridade contra Solver 100k.",
        "FlatNumba domina em alto budget.",
        "Trade-off claro: qualidade vs latência.",
        "Forte vantagem do primeiro a jogar.",
    ], size=10, spacing=2)

    footer(s, 16)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 17 — DIFICULDADES & SOLUÇÕES
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)
    header(s, "Dificuldades & Soluções", chip="05")

    pairs = [
        ("Tree reuse com flat arrays",
         "BFS compaction re-indexa pai/filhos antes do kernel.",
         TEAL,
         "01"),
        ("Data leakage não-óbvio (state-level)",
         "Auditoria + 3 splits comparados  →  split agrupado por estado.",
         CORAL,
         "02"),
        ("Classes desequilibradas (89 % drops)",
         "Oversampling controlado de proven positions (× 4).",
         AMBER,
         "03"),
        ("Performance Python insuficiente",
         "Kernels Numba JIT em arrays planos  →  ~44× speedup.",
         FOREST,
         "04"),
    ]
    py0 = 1.25
    ph = 0.85
    pgap = 0.15
    for i, (problem, solution, col, num) in enumerate(pairs):
        y = py0 + i * (ph + pgap)
        add_rect(s, 0.5, y, 9.0, ph, fill=ICE)
        add_rect(s, 0.5, y, 0.15, ph, fill=col)
        add_text(s, 0.75, y + 0.10, 0.55, 0.65, num,
                 size=22, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE, font=HEADER_FONT)
        add_text(s, 1.40, y + 0.10, 8.0, 0.35, problem,
                 size=13, bold=True, color=DARK)
        add_text(s, 1.40, y + 0.45, 8.0, 0.35, "→  " + solution,
                 size=11, color=SLATE, italic=True)

    footer(s, 17)

    # ════════════════════════════════════════════════════════════════════════
    # SLIDE 18 — CONCLUSÃO
    # ════════════════════════════════════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, NAVY)

    add_rect(s, 0.0, 0.0, 0.30, PAGE_H, fill=TEAL)
    # Decorative discs
    for i in range(5):
        add_circle(s, 8.30 + i * 0.30, 0.40, 0.22,
                   fill=[AMBER, WHITE, TEAL_LT, AMBER, WHITE][i])

    add_text(s, 1.0, 0.40, 8.0, 0.55,
             "Conclusão", size=36, bold=True, color=WHITE,
             font=HEADER_FONT)
    add_rect(s, 1.0, 1.00, 1.8, 0.04, fill=AMBER)

    # Three columns
    blocks = [
        ("Requisitos do guião",
         [
             "3 cenários (HvH/HvC/CvC)",
             "MCTS com UCT",
             "Variantes além do standard",
             "ID3 from scratch",
             "Iris + dataset PopOut",
         ],
         AMBER),
        ("Contribuições técnicas",
         [
             "MCTS-Solver com proof propagation",
             "Tree reuse com BFS compaction",
             "Numba: ~220 k iter/s  ·  ~44×",
             "~2 200 linhas de testes",
             "Pipeline reprodutível",
         ],
         TEAL),
        ("Rigor & honestidade",
         [
             "3 splits comparados",
             "Híbrido declarado explicitamente",
             "Limitações documentadas",
             "Baselines reportadas",
             "Notebook auto-contido",
         ],
         FOREST),
    ]
    for i, (title, items, col) in enumerate(blocks):
        x = 1.0 + i * 2.95
        add_text(s, x, 1.20, 2.8, 0.35, title,
                 size=13, bold=True, color=col, font=HEADER_FONT)
        for j, item in enumerate(items):
            y = 1.60 + j * 0.34
            # check
            add_text(s, x, y, 0.18, 0.30, "✓",
                     size=12, bold=True, color=col,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(s, x + 0.22, y, 2.6, 0.30, item,
                     size=11, color=ICE, anchor=MSO_ANCHOR.MIDDLE)

    # Future work bar
    add_rect(s, 1.0, 4.40, 8.0, 0.55, fill=DEEP)
    add_text(s, 1.0, 4.40, 8.0, 0.27,
             "TRABALHO FUTURO", size=10, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER)
    add_text(s, 1.0, 4.65, 8.0, 0.30,
             "Zobrist hashing + transposition tables   ·   AMAF / RAVE   ·   Opening books",
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Thanks + link
    add_text(s, 1.0, 5.10, 8.0, 0.25,
             "Obrigado  ·  github.com/Difl4/popout-ai",
             size=11, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)

    # ─── Save ───────────────────────────────────────────────────────────────
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"  {len(prs.slides)} slides  ·  {OUT_PATH.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
