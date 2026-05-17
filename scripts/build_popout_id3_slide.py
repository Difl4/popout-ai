"""Build a single standalone slide about ID3 on PopOut (the real ML task).

Output: docs/popout_id3_slide.pptx — one slide (16:9), matching the visual
language of the Numba / Solver+Reuse / Iris slides.

All numbers come directly from the v4_75k_100k results saved by the project:
  data/generated/v4_75k_100k/results/{accuracy_results,winrate_results}.json

  Dataset
    - 75 000 self-play games
    - Oracle: ReuseFlatNumbaSolverMCTS at 100 k iterations
    - 2,210,901 train + 552,725 test rows under random splits;
      2,095,174 + 668,452 under group split.
    - Mirroring + tiered blunder rates (5 / 25 / 50 %)

  Three-split methodology (test accuracy):
    Split     overlap   feat       raw
    BUGGY     98.5 %    0.868     0.870
    RANDOM    92.3 %    0.819     0.823
    GROUPED    0.0 %    0.567     0.584   ← canonical

  Tactical accuracy on GROUPED test (the hardest sub-positions):
    Metric            feat     raw    support
    opp_wins_next=1   0.372   0.392    89 345
    pop moves         0.290   0.337    34 737

  Win-rate vs FlatNumbaSolverMCTS @ 100 k iter (100 games, alternating colours):
    Agent                          Win-rate
    ID3 (features) hybrid          53 %
    ID3 Raw         hybrid         54 %
    ID3 (features) pure            36 %       ← without _forced_move
    ID3 Raw         pure           29 %       ← without _forced_move

Usage:
    python scripts/build_popout_id3_slide.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs/popout_id3_slide.pptx"

# ─── Palette ──────────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x21, 0x29, 0x5C)
DEEP     = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT  = RGBColor(0x5A, 0xA4, 0xB8)
ICE      = RGBColor(0xE6, 0xF1, 0xF5)
ICE_DK   = RGBColor(0xC4, 0xDB, 0xE3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x12, 0x1E, 0x3F)
SLATE    = RGBColor(0x33, 0x44, 0x55)
GRAY     = RGBColor(0x6B, 0x7A, 0x88)
AMBER    = RGBColor(0xF4, 0xA3, 0x41)
FOREST   = RGBColor(0x2E, 0x8B, 0x57)
CORAL    = RGBColor(0xE5, 0x71, 0x4E)
RED      = RGBColor(0xC1, 0x3F, 0x2F)

HEADER_FONT = "Cambria"
BODY_FONT   = "Calibri"


# ─── Primitives ──────────────────────────────────────────────────────────────
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def add_bullets(slide, x, y, w, h, items, *, size=10, color=SLATE,
                bullet_color=TEAL, spacing=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        b = p.add_run()
        b.text = "•  "
        b.font.name = BODY_FONT
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=None, line=None,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


# ─── Build slide ─────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # ── Header ───────────────────────────────────────────────────────────────
    add_rect(s, 0.5, 0.25, 0.65, 0.33, fill=TEAL)
    add_text(s, 0.5, 0.25, 0.65, 0.33, "06",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.30, 0.21, 8.2, 0.45,
             "ID3 on PopOut — Honest Generalisation",
             size=22, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 1.30, 0.62, 8.2, 0.28,
             "3-split methodology exposed a data-leakage bug; true generalisation "
             "is 56-58 %, not the 87 % a naïve split suggests",
             size=10, italic=True, color=GRAY)
    add_rect(s, 0.5, 0.97, 9.0, 0.015, fill=ICE_DK)

    # ── Two cards: Dataset | Methodology ─────────────────────────────────────
    card_y = 1.10
    card_h = 2.25

    # Left — Dataset
    add_rect(s, 0.5, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 0.5, card_y, 4.40, 0.10, fill=TEAL)
    add_text(s, 0.65, card_y + 0.16, 4.10, 0.32,
             "Dataset (v4_75k_100k)", size=15, bold=True, color=TEAL,
             font=HEADER_FONT)
    add_text(s, 0.65, card_y + 0.48, 4.10, 0.24,
             "75 000 self-play games · oracle = ReuseFlat Solver @ 100k",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 0.65, card_y + 0.78, 4.10, 1.40, [
        "2.76 M (state, best_move) pairs from full self-play games "
        "labelled by the oracle.",
        "Forced 2-move openings cycling 7 × 7 = 49 column combos for "
        "coverage.",
        "Three blunder tiers (5 / 25 / 50 %) for diversity in opponent "
        "play.",
        "Horizontal mirroring doubles useful samples; 39 % of positions "
        "carry the is_proven flag.",
    ], size=10, bullet_color=TEAL)

    # Right — Methodology (3-split)
    add_rect(s, 5.10, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 5.10, card_y, 4.40, 0.10, fill=AMBER)
    add_text(s, 5.25, card_y + 0.16, 4.10, 0.32,
             "3-split methodology", size=15, bold=True, color=DEEP,
             font=HEADER_FONT)
    add_text(s, 5.25, card_y + 0.48, 4.10, 0.24,
             "Exposed our own data leakage — reported honestly",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 5.25, card_y + 0.78, 4.10, 1.40, [
        "BUGGY: reset_index bug → 98.5 % literal row overlap "
        "(memorisation, not generalisation).",
        "RANDOM: correct row-level random 80/20 → 92 % overlap "
        "from intrinsic state duplication.",
        "GROUPED: states partitioned 80/20 before mapping back to "
        "rows → 0 % overlap, true generalisation.",
        "Reported all three side-by-side in the notebook.",
    ], size=10, bullet_color=AMBER)

    # ── Middle: test-accuracy comparison chart + callout ─────────────────────
    add_text(s, 0.50, 3.45, 6.40, 0.25,
             "Test accuracy across the three splits",
             size=11, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 0.50, 3.68, 6.40, 0.18,
             "Features model | Raw model — same trend: GROUPED is the honest measure",
             size=8, italic=True, color=GRAY)

    # Grouped horizontal bars (3 splits × 2 models)
    bars = [
        ("BUGGY",   0.868, 0.870, CORAL),
        ("RANDOM",  0.819, 0.823, AMBER),
        ("GROUPED", 0.567, 0.584, FOREST),
    ]
    chart_y0 = 3.92
    label_x  = 0.50
    label_w  = 1.30
    chart_x  = label_x + label_w + 0.05
    chart_w  = 3.40
    pair_h   = 0.30     # per split, two stacked thin bars
    bar_h    = 0.13
    bar_gap  = 0.04
    group_gap = 0.05

    for i, (split, feat, raw, accent) in enumerate(bars):
        y_top = chart_y0 + i * (pair_h + group_gap)
        # split label centred between the two bars
        add_text(s, label_x, y_top, label_w, pair_h, split,
                 size=11, bold=True, color=accent,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # feat bar
        add_rect(s, chart_x, y_top, chart_w * feat, bar_h, fill=accent)
        add_text(s, chart_x + chart_w * feat + 0.05, y_top, 0.85, bar_h,
                 f"feat  {feat:.3f}",
                 size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        # raw bar (slightly lighter shade via transparency would be ideal;
        # instead we use the same colour with a horizontal offset rectangle)
        y_raw = y_top + bar_h + bar_gap
        add_rect(s, chart_x, y_raw, chart_w * raw, bar_h, fill=accent)
        # add a horizontal stripe to distinguish "raw"
        add_rect(s, chart_x, y_raw, chart_w * raw, 0.04, fill=WHITE)
        add_text(s, chart_x + chart_w * raw + 0.05, y_raw, 0.85, bar_h,
                 f"raw   {raw:.3f}",
                 size=9, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    # ── Right: TRUE GENERALISATION callout ───────────────────────────────────
    add_rect(s, 7.10, 3.45, 2.40, 1.50, fill=NAVY)
    add_text(s, 7.10, 3.53, 2.40, 0.24,
             "TRUE GENERALISATION", size=9, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER)
    add_text(s, 7.10, 3.82, 2.40, 0.55,
             "~ 58 %",
             size=30, bold=True, color=WHITE, font=HEADER_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.10, 4.42, 2.40, 0.24,
             "vs 87 % under leakage",
             size=10, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)
    add_text(s, 7.10, 4.68, 2.40, 0.24,
             "random baseline:  7 %  (1/14)",
             size=8, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)

    # ── Bottom strip: tactical + hybrid honesty ──────────────────────────────
    strip_y = 5.05
    strip_h = 0.45

    # Left — tactical accuracy
    add_rect(s, 0.5, strip_y, 5.70, strip_h, fill=ICE)
    add_rect(s, 0.5, strip_y, 0.10, strip_h, fill=DEEP)
    add_text(s, 0.70, strip_y + 0.04, 5.40, 0.20,
             "Tactical accuracy on GROUPED test",
             size=9, bold=True, color=NAVY)
    add_text(s, 0.70, strip_y + 0.23, 5.40, 0.20,
             "opp_wins_next=1: 37 % (feat)  ·  39 % (raw)   |   "
             "pop moves: 29 % (feat)  ·  34 % (raw)",
             size=8, color=SLATE, italic=True)

    # Right — hybrid honesty
    add_rect(s, 6.30, strip_y, 3.20, strip_h, fill=ICE)
    add_rect(s, 6.30, strip_y, 0.10, strip_h, fill=CORAL)
    add_text(s, 6.50, strip_y + 0.04, 2.95, 0.20,
             "Hybrid disclosure", size=9, bold=True, color=NAVY)
    add_text(s, 6.50, strip_y + 0.23, 2.95, 0.20,
             "Pure ID3 vs Solver 100k: 29-36 %  ·  hybrid: 53-54 %",
             size=8, color=SLATE, italic=True)

    # ── Save ─────────────────────────────────────────────────────────────────
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"  1 slide · {OUT_PATH.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
