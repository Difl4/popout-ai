"""Build a single standalone slide about Numba + FlatNumba.

Output: docs/numba_slide.pptx — one slide (16:9), ready to insert into any deck.
Uses the same "Ocean Gradient" design system as the main presentation for
visual continuity.

Usage:
    python scripts/build_numba_slide.py
"""
from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs/numba_slide.pptx"

# ─── Palette (matches main presentation) ──────────────────────────────────────
NAVY     = RGBColor(0x21, 0x29, 0x5C)
DEEP     = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT  = RGBColor(0x5A, 0xA4, 0xB8)
ICE      = RGBColor(0xE6, 0xF1, 0xF5)
ICE_DK   = RGBColor(0xC4, 0xDB, 0xE3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x12, 0x1E, 0x3F)
SLATE    = RGBColor(0x33, 0x44, 0x55)
GRAY     = RGBColor(0x6B, 0x7A, 0x88)
AMBER    = RGBColor(0xF4, 0xA3, 0x41)
FOREST   = RGBColor(0x2E, 0x8B, 0x57)

HEADER_FONT = "Cambria"
BODY_FONT   = "Calibri"
MONO_FONT   = "Consolas"


# ─── Primitives ──────────────────────────────────────────────────────────────
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def add_bullets(slide, x, y, w, h, items, *, size=10, color=SLATE,
                bullet_color=TEAL, spacing=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        b = p.add_run()
        b.text = "•  "
        b.font.name = BODY_FONT
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=None, line=None,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


# ─── Build slide ─────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # ── Header ───────────────────────────────────────────────────────────────
    add_rect(s, 0.5, 0.25, 0.65, 0.33, fill=TEAL)
    add_text(s, 0.5, 0.25, 0.65, 0.33, "03",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.30, 0.21, 8.2, 0.45,
             "Numba JIT + FlatNumba",
             size=22, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 1.30, 0.62, 8.2, 0.28,
             "From Python to native speed — ~44× more state visits per second, "
             "without leaving Python",
             size=10, italic=True, color=GRAY)
    add_rect(s, 0.5, 0.97, 9.0, 0.015, fill=ICE_DK)

    # ── Two cards: Numba JIT | FlatNumba ─────────────────────────────────────
    card_y = 1.10
    card_h = 2.25

    # Left — Numba JIT
    add_rect(s, 0.5, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 0.5, card_y, 4.40, 0.10, fill=TEAL)
    add_text(s, 0.65, card_y + 0.16, 4.10, 0.32,
             "Numba JIT", size=15, bold=True, color=TEAL,
             font=HEADER_FONT)
    add_text(s, 0.65, card_y + 0.48, 4.10, 0.24,
             "@njit compiles Python → native via LLVM",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 0.65, card_y + 0.78, 4.10, 1.40, [
        "Decorator on the four MCTS hot loops "
        "— select · expand · simulate · backprop.",
        "@njit(cache=True) persists compiled bytecode across runs; "
        "explicit warmup() pays the JIT cost up front.",
        "All operations on primitive types (int / float / numpy arrays) — "
        "no Python object access inside the kernel.",
        "Chosen over Cython / C extensions: same ~95% of the speed-up, "
        "zero extra build chain.",
    ], size=10)

    # Right — FlatNumba
    add_rect(s, 5.10, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 5.10, card_y, 4.40, 0.10, fill=AMBER)
    add_text(s, 5.25, card_y + 0.16, 4.10, 0.32,
             "FlatNumba", size=15, bold=True, color=DEEP,
             font=HEADER_FONT)
    add_text(s, 5.25, card_y + 0.48, 4.10, 0.24,
             "Whole tree as pre-allocated NumPy arrays",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 5.25, card_y + 0.78, 4.10, 1.40, [
        "Nodes are integer indices into fixed-size arrays "
        "(visits · value · parent · children · status).",
        "Memory layout is contiguous → cache-friendly; "
        "no GC pressure, no Python object dispatch.",
        "Tree reuse: BFS compaction re-indexes the surviving "
        "subtree to slots 0..k−1 between turns.",
        "Visit budget compounds — each turn inherits ≈12% of "
        "the previous turn's search effort.",
    ], size=10, bullet_color=AMBER)

    # ── Performance bar chart (log-scale, units: iter/s) ─────────────────────
    # Numbers measured on the project's own benchmark (Technical_Documentation
    # notebook + fresh re-run on the dev machine). Slide scope is Numba +
    # FlatNumba only; tree-reuse / Solver variants live in a separate slide
    # because reuse does not change raw iter/s.
    bars = [
        ("StandardUCT",     "Python baseline",   5_000, TEAL_LT),
        ("NumbaMCTS",       "Objects + JIT",    30_000, TEAL),
        ("FlatNumbaMCTS",   "Arrays + JIT",    220_000, DEEP),
    ]
    chart_y0  = 3.55
    label_x   = 0.50
    label_w   = 1.75
    desc_x    = label_x + label_w + 0.05
    desc_w    = 1.40
    chart_x   = desc_x + desc_w + 0.05
    chart_w   = 2.60
    bar_h     = 0.30
    bar_gap   = 0.10
    bar_max   = max(b[2] for b in bars)

    for i, (name, desc, val, col) in enumerate(bars):
        y = chart_y0 + i * (bar_h + bar_gap)
        w = max(0.10,
                math.log10(val + 1) / math.log10(bar_max + 1) * chart_w)
        add_text(s, label_x, y, label_w, bar_h, name,
                 size=10, bold=True, color=DARK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, desc_x, y, desc_w, bar_h, desc,
                 size=8, italic=True, color=GRAY,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, chart_x, y, w, bar_h, fill=col)
        if val >= 1_000_000:
            v_str = f"{val // 1_000_000} M / s"
        elif val >= 1000:
            v_str = f"{val // 1000} k / s"
        else:
            v_str = f"{val} / s"
        add_text(s, chart_x + w + 0.05, y, 0.95, bar_h, v_str,
                 size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)

    # Chart caption + raw/effective disclaimer
    add_text(s, label_x, chart_y0 + 4 * (bar_h + bar_gap) + 0.04,
             chart_w + chart_x + 0.5, 0.20,
             "log scale  ·  state visits per second of wall-clock",
             size=8, color=GRAY, italic=True)

    # ── Speedup callout (right side) ─────────────────────────────────────────
    add_rect(s, 7.85, 3.55, 1.65, 1.20, fill=NAVY)
    add_text(s, 7.85, 3.60, 1.65, 0.25,
             "SPEEDUP", size=9, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER)
    add_text(s, 7.85, 3.85, 1.65, 0.60,
             "≈ 44×",
             size=32, bold=True, color=WHITE, font=HEADER_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.85, 4.45, 1.65, 0.25,
             "vs pure Python", size=8, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)

    # ── Bottom strip: raw vs effective + real-world impact ───────────────────
    strip_y = 4.95
    strip_h = 0.50

    # Left half — speedup composition
    add_rect(s, 0.5, strip_y, 5.70, strip_h, fill=ICE)
    add_rect(s, 0.5, strip_y, 0.10, strip_h, fill=DEEP)
    add_text(s, 0.70, strip_y + 0.04, 5.40, 0.22,
             "Where the ~44× comes from",
             size=10, bold=True, color=NAVY)
    add_text(s, 0.70, strip_y + 0.25, 5.40, 0.22,
             "Numba JIT alone: ~6×  (5 k → 30 k)   ·   "
             "Flat-array layout: another ~7×  (30 k → 220 k).",
             size=9, color=SLATE, italic=True)

    # Right half — practical impact
    add_rect(s, 6.30, strip_y, 3.20, strip_h, fill=ICE)
    add_rect(s, 6.30, strip_y, 0.10, strip_h, fill=FOREST)
    add_text(s, 6.50, strip_y + 0.04, 2.95, 0.22,
             "Practical impact", size=10, bold=True, color=NAVY)
    add_text(s, 6.50, strip_y + 0.25, 2.95, 0.22,
             "Competitive moves in ~50 ms · real-time GUI play.",
             size=9, color=SLATE, italic=True)

    # ── Save ─────────────────────────────────────────────────────────────────
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"  1 slide · {OUT_PATH.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
