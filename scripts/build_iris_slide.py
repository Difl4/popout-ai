"""Build a single standalone slide about ID3 on Iris (warm-up validation).

Output: docs/iris_slide.pptx — one slide (16:9), matching the visual language
of the Numba and Solver+Reuse slides.

All numbers are extracted directly from the project's own evaluation:
notebooks/ID3_Decision_Tree.ipynb cells 8, 12, 14.

  Methodology
    - Quantile discretisation, q = 3 bins per feature
    - 5 × 10-fold cross-validation  =  50 independent train/test runs
    - Discretiser fit inside each fold on the training partition only
      (no leakage from test into train)

  Accuracy across 50 folds
    Mean : 95.07 %
    Std  :  5.63 %
    Min  : 80.00 %
    Max  : 100.00 %

  Per-class metrics (support = 250 each)
    Class            Precision   Recall   F1
    Iris-setosa        0.996     0.976   0.986
    Iris-versicolor    0.938     0.912   0.925
    Iris-virginica     0.920     0.964   0.941

  Feature importance (normalised split count)
    sepalwidth   33.33 %
    sepallength  33.33 %
    petalwidth   22.22 %
    petallength  11.11 %

Usage:
    python scripts/build_iris_slide.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ─── Paths ────────────────────────────────────────────────────────────────────
ROOT     = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "docs/iris_slide.pptx"

# ─── Palette (matches the other standalone slides) ────────────────────────────
NAVY     = RGBColor(0x21, 0x29, 0x5C)
DEEP     = RGBColor(0x06, 0x5A, 0x82)
TEAL     = RGBColor(0x1C, 0x72, 0x93)
TEAL_LT  = RGBColor(0x5A, 0xA4, 0xB8)
ICE      = RGBColor(0xE6, 0xF1, 0xF5)
ICE_DK   = RGBColor(0xC4, 0xDB, 0xE3)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x12, 0x1E, 0x3F)
SLATE    = RGBColor(0x33, 0x44, 0x55)
GRAY     = RGBColor(0x6B, 0x7A, 0x88)
AMBER    = RGBColor(0xF4, 0xA3, 0x41)
FOREST   = RGBColor(0x2E, 0x8B, 0x57)
ROSE     = RGBColor(0xC8, 0x59, 0x7A)

HEADER_FONT = "Cambria"
BODY_FONT   = "Calibri"


# ─── Primitives (identical to the other build scripts) ───────────────────────
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=BODY_FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def add_bullets(slide, x, y, w, h, items, *, size=10, color=SLATE,
                bullet_color=TEAL, spacing=4):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_top = Inches(0.0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        b = p.add_run()
        b.text = "•  "
        b.font.name = BODY_FONT
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = BODY_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color


def add_rect(slide, x, y, w, h, *, fill=None, line=None,
             shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


# ─── Build slide ─────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    set_bg(s, WHITE)

    # ── Header ───────────────────────────────────────────────────────────────
    add_rect(s, 0.5, 0.25, 0.65, 0.33, fill=TEAL)
    add_text(s, 0.5, 0.25, 0.65, 0.33, "05",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 1.30, 0.21, 8.2, 0.45,
             "ID3 on Iris — Warm-up Validation",
             size=22, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 1.30, 0.62, 8.2, 0.28,
             "Validating the from-scratch ID3 on the canonical benchmark "
             "before the PopOut dataset",
             size=10, italic=True, color=GRAY)
    add_rect(s, 0.5, 0.97, 9.0, 0.015, fill=ICE_DK)

    # ── Two cards: Methodology | Per-class metrics ───────────────────────────
    card_y = 1.10
    card_h = 2.25

    # Left — Methodology
    add_rect(s, 0.5, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 0.5, card_y, 4.40, 0.10, fill=TEAL)
    add_text(s, 0.65, card_y + 0.16, 4.10, 0.32,
             "Methodology", size=15, bold=True, color=TEAL,
             font=HEADER_FONT)
    add_text(s, 0.65, card_y + 0.48, 4.10, 0.24,
             "Quantile discretisation + repeated K-fold CV — leakage-safe",
             size=9, italic=True, color=GRAY)
    add_bullets(s, 0.65, card_y + 0.78, 4.10, 1.40, [
        "Continuous features → 3 quantile bins per attribute "
        "(very_low / low / medium).",
        "5 × 10-fold stratified CV → 50 independent train/test "
        "evaluations.",
        "Bin edges fit on training fold only — test partition never "
        "influences the splits.",
        "ID3 from src/decision_tree/id3/learner.py — same algorithm "
        "used downstream for PopOut.",
    ], size=10, bullet_color=TEAL)

    # Right — Per-class metrics (mini table)
    add_rect(s, 5.10, card_y, 4.40, card_h, fill=ICE)
    add_rect(s, 5.10, card_y, 4.40, 0.10, fill=AMBER)
    add_text(s, 5.25, card_y + 0.16, 4.10, 0.32,
             "Per-class metrics", size=15, bold=True, color=DEEP,
             font=HEADER_FONT)
    add_text(s, 5.25, card_y + 0.48, 4.10, 0.24,
             "Support = 250 samples per class  ·  aggregated over 50 folds",
             size=9, italic=True, color=GRAY)

    # Mini table — header row
    tx0 = 5.30
    col_w_lab = 1.65
    col_w_num = 0.78
    th_y      = card_y + 0.82
    th_h      = 0.20
    add_text(s, tx0,                              th_y, col_w_lab, th_h,
             "Class", size=9, bold=True, color=NAVY)
    add_text(s, tx0 + col_w_lab,                  th_y, col_w_num, th_h,
             "Prec.", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, tx0 + col_w_lab + col_w_num,      th_y, col_w_num, th_h,
             "Recall", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, tx0 + col_w_lab + 2 * col_w_num,  th_y, col_w_num, th_h,
             "F1", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(s, tx0, th_y + 0.22, col_w_lab + 3 * col_w_num, 0.015,
             fill=ICE_DK)

    classes = [
        ("Iris-setosa",     "0.996", "0.976", "0.986", FOREST),
        ("Iris-versicolor", "0.938", "0.912", "0.925", DEEP),
        ("Iris-virginica",  "0.920", "0.964", "0.941", ROSE),
    ]
    row_h = 0.28
    for i, (name, p, r, f1, col) in enumerate(classes):
        y = th_y + 0.30 + i * row_h
        add_rect(s, tx0 - 0.05, y, 0.08, row_h - 0.04, fill=col)
        add_text(s, tx0 + 0.05, y, col_w_lab - 0.05, row_h,
                 name, size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx0 + col_w_lab,             y, col_w_num, row_h,
                 p, size=10, color=SLATE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx0 + col_w_lab + col_w_num, y, col_w_num, row_h,
                 r, size=10, color=SLATE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, tx0 + col_w_lab + 2 * col_w_num, y, col_w_num, row_h,
                 f1, size=10, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── Middle: F1 chart + big accuracy callout ──────────────────────────────
    add_text(s, 0.50, 3.45, 6.40, 0.25,
             "F1-score per class — 50-fold aggregated",
             size=11, bold=True, color=NAVY, font=HEADER_FONT)
    add_text(s, 0.50, 3.68, 6.40, 0.18,
             "Bars filled to F1 score; thin grey background bar shows 0–1 range",
             size=8, italic=True, color=GRAY)

    bars = [
        ("Iris-setosa",     0.986, FOREST),
        ("Iris-versicolor", 0.925, DEEP),
        ("Iris-virginica",  0.941, ROSE),
    ]
    chart_y0 = 3.95
    label_x  = 0.50
    label_w  = 1.65
    chart_x  = label_x + label_w + 0.05
    chart_w  = 3.20
    bar_h    = 0.26
    bar_gap  = 0.10

    for i, (name, val, col) in enumerate(bars):
        y = chart_y0 + i * (bar_h + bar_gap)
        # background "0..1" bar
        add_rect(s, chart_x, y, chart_w, bar_h, fill=ICE)
        # foreground bar (length proportional to F1)
        add_rect(s, chart_x, y, chart_w * val, bar_h, fill=col)
        # label on left
        add_text(s, label_x, y, label_w, bar_h, name,
                 size=10, bold=True, color=DARK,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        # value on right
        add_text(s, chart_x + chart_w + 0.05, y, 0.65, bar_h,
                 f"{val:.3f}",
                 size=10, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)

    # ── Right: accuracy callout (aligned with chart) ─────────────────────────
    add_rect(s, 7.10, 3.45, 2.40, 1.50, fill=NAVY)
    add_text(s, 7.10, 3.53, 2.40, 0.24,
             "OVERALL ACCURACY", size=9, bold=True, color=AMBER,
             align=PP_ALIGN.CENTER)
    add_text(s, 7.10, 3.82, 2.40, 0.55,
             "95.07 %",
             size=28, bold=True, color=WHITE, font=HEADER_FONT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 7.10, 4.40, 2.40, 0.24,
             "± 5.63 %  (50 folds)",
             size=10, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)
    add_text(s, 7.10, 4.66, 2.40, 0.24,
             "min 80 %   ·   max 100 %",
             size=8, italic=True, color=TEAL_LT,
             align=PP_ALIGN.CENTER)

    # ── Bottom: feature importance + reproducibility note ────────────────────
    strip_y = 5.05
    strip_h = 0.45

    # Left — feature importance
    add_rect(s, 0.5, strip_y, 5.70, strip_h, fill=ICE)
    add_rect(s, 0.5, strip_y, 0.10, strip_h, fill=DEEP)
    add_text(s, 0.70, strip_y + 0.04, 5.40, 0.20,
             "Feature importance (normalised split count)",
             size=9, bold=True, color=NAVY)
    add_text(s, 0.70, strip_y + 0.23, 5.40, 0.20,
             "sepal-width 33 %  ·  sepal-length 33 %  ·  "
             "petal-width 22 %  ·  petal-length 11 %",
             size=8, color=SLATE, italic=True)

    # Right — reproducibility note
    add_rect(s, 6.30, strip_y, 3.20, strip_h, fill=ICE)
    add_rect(s, 6.30, strip_y, 0.10, strip_h, fill=FOREST)
    add_text(s, 6.50, strip_y + 0.04, 2.95, 0.20,
             "Reproducibility", size=9, bold=True, color=NAVY)
    add_text(s, 6.50, strip_y + 0.23, 2.95, 0.20,
             "Same ID3Classifier reused for PopOut downstream.",
             size=8, color=SLATE, italic=True)

    # ── Save ─────────────────────────────────────────────────────────────────
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"  1 slide · {OUT_PATH.stat().st_size // 1024} KB")


if __name__ == "__main__":
    build()
